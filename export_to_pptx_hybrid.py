#!/usr/bin/env python3
"""
Hybrid PPTX export — flow diagrams stay as images, everything else is editable.

  Slide 1  Cover         — full image (orb gradient not reproducible)
  Slide 2  North Star    — pipeline region as image; label, title, sub, 3 stat cards editable
  Slide 3  Problem       — example + flow as image; label, title, sub, 4 stat cards editable
  Slide 4  Hypothesis    — example + flow as image; label, title, sub, 4 benefit cards editable
  Slide 5  Solution      — top flow section as image; label, title, 5 build cards editable
  Slide 6  Roadmap       — roadmap phase cards as image; label, title, subtitle editable
"""

import asyncio
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from playwright.async_api import async_playwright

BASE  = Path('/Users/jennywang/Documents/Claude/New Fulfillment Workflow')
SHOTS = BASE / '_hybrid'
OUT   = BASE / 'New Fulfillment Flow Vision — Hybrid.pptx'
SHOTS.mkdir(exist_ok=True)

SCALE = 10.0 / 960   # px → inches (960px slide = 10in)
TOP   = 68            # page offset to slide top

def i(px):  return Inches(px * SCALE)
def pt(px): return Pt(px * 0.75)        # CSS px → PPTX points

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# ── Design-system colors ──────────────────────────────────────────────────────
PRIMARY     = rgb('#2945f0')
PRIMARY_800 = rgb('#142592')
PRIMARY_LT  = rgb('#eceefe')
BODY        = rgb('#020202')
DIMMED      = rgb('#687382')
BORDER      = rgb('#c2c7ce')
SURF        = rgb('#f2f3f8')
DARK_NAVY   = rgb('#1d2024')
TEAL        = rgb('#0d5c63')
TEAL_LT     = rgb('#e0f2f4')
WHITE       = rgb('#ffffff')
RED         = rgb('#cc3300')
ORANGE      = rgb('#ea580c')
AMBER       = rgb('#b45309')
AMBER_BD    = rgb('#f59e0b')
SUCCESS     = rgb('#2e6d57')
GRAY_MID    = rgb('#6b7280')
GRAY_LIGHT  = rgb('#9ca3af')


# ── Shape helpers ─────────────────────────────────────────────────────────────

def _no_line(shp):  shp.line.fill.background()
def _solid(shp, c): shp.fill.solid(); shp.fill.fore_color.rgb = c
def _no_fill(shp):  shp.fill.background()

def _no_inset(tf):
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        for attr in ('lIns', 'tIns', 'rIns', 'bIns'):
            bp.set(attr, '0')

def add_rect(slide, x, y, w, h, fill=None, line_color=None, lw=0.75):
    shp = slide.shapes.add_shape(1, i(x), i(y), i(w), i(h))
    if fill:    _solid(shp, fill)
    else:       _no_fill(shp)
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(lw)
    else:
        _no_line(shp)
    return shp

def add_text(slide, text, x, y, w, h,
             size=9, bold=False, italic=False, color=BODY,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    for idx, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text   = line
        run.font.name   = 'Geist'
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
    _no_inset(tf)
    return txb

def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(str(path), i(x), i(y), i(w), i(h))


# ── JS helper to get element bounding boxes ──────────────────────────────────

JS_BOXES = """(selectors) => {
    const TOP = 68;
    const out = {};
    for (const [k, sel] of Object.entries(selectors)) {
        out[k] = [...document.querySelectorAll(sel)].map(el => {
            const r = el.getBoundingClientRect();
            return { x: r.left, y: r.top - TOP, w: r.width, h: r.height,
                     text: el.innerText.trim() };
        });
    }
    return out;
}"""

async def boxes(page, path, sels):
    await page.goto(f'file://{path}')
    await page.wait_for_timeout(1200)
    return await page.evaluate(JS_BOXES, sels)

async def clip_shot(page, path, y0, y1, out_path):
    await page.goto(f'file://{path}')
    await page.wait_for_timeout(1200)
    h = max(1, int(y1 - y0))
    await page.screenshot(
        path=str(out_path),
        clip={'x': 0, 'y': int(TOP + y0), 'width': 960, 'height': h}
    )
    return h


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_bg(slide):
    """Add a white background rectangle covering the full slide."""
    add_rect(slide, 0, 0, 960, 540, fill=WHITE)

def _header_lbl(slide, b):
    """Place the slide label from the bounding box dict (first element)."""
    if b and len(b) > 0:
        el = b[0]
        add_text(slide, el['text'], el['x'], el['y'], el['w'] + 20, el['h'] + 6,
                 size=7.5, color=PRIMARY)

def _header_h1(slide, b, max_width=912):
    if b and len(b) > 0:
        el = b[0]
        add_text(slide, el['text'], el['x'], el['y'], max_width, el['h'] + 20,
                 size=13.5, bold=True, color=BODY)

def _header_sub(slide, b, lborder=None):
    """Place subtitle text with optional left border bar."""
    if not b or len(b) == 0: return
    el = b[0]
    if lborder:
        # 3px colored left border bar
        add_rect(slide, el['x'], el['y'], 3, el['h'] + 4, fill=lborder)
        add_text(slide, el['text'], el['x'] + 7, el['y'], el['w'], el['h'] + 8,
                 size=7.5, italic=False, color=DIMMED, wrap=True)
    else:
        add_text(slide, el['text'], el['x'], el['y'], el['w'], el['h'] + 8,
                 size=7.5, color=DIMMED, wrap=True)


# ── SLIDE 1: Cover — full image ───────────────────────────────────────────────

async def build_slide1(prs, page, layout):
    print('  Slide 1 — Cover (full image) …')
    path = BASE / 'slide_01_cover.html'
    out  = SHOTS / 'slide_01.png'
    await page.goto(f'file://{path}')
    await page.wait_for_timeout(1200)
    await page.screenshot(path=str(out),
                          clip={'x': 0, 'y': TOP, 'width': 960, 'height': 540})
    s = prs.slides.add_slide(layout)
    add_image(s, out, 0, 0, 960, 540)
    return s


# ── SLIDE 2: ROI North Star ───────────────────────────────────────────────────

async def build_slide2(prs, page, layout):
    print('  Slide 2 — ROI North Star …')
    path  = BASE / 'slide_02_north_star.html'
    b = await boxes(page, path, {
        'lbl':   '.slide-label',
        'h1':    '.h1',
        'sub':   '.sub',
        'stats': '.stat-card',
    })

    # Diagram region: bottom of .sub → top of first .stat-card
    sub_bot = b['sub'][0]['y'] + b['sub'][0]['h'] if b['sub'] else 80
    stats_top = b['stats'][0]['y'] if b['stats'] else 310

    out = SHOTS / 'slide_02_diag.png'
    diag_h = await clip_shot(page, path, sub_bot + 6, stats_top - 6, out)

    s = prs.slides.add_slide(layout)
    _slide_bg(s)

    # Header
    _header_lbl(s, b['lbl'])
    _header_h1(s, b['h1'])
    add_text(s, b['sub'][0]['text'],
             b['sub'][0]['x'], b['sub'][0]['y'],
             b['sub'][0]['w'], b['sub'][0]['h'] + 8,
             size=9, color=DIMMED, wrap=True)

    # Diagram image
    add_image(s, out, 0, sub_bot + 6, 960, diag_h)

    # Stat cards — editable
    stat_data = [
        ('2 min → 50 sec', 'processing time per request from Phases 1–3 — AI extraction reducing manual logging and auth work at scale', PRIMARY, SURF),
        ('32M requests', 'processed in 2025 — every minute of fulfillment time eliminated across this volume = significant labor cost out', PRIMARY, SURF),
        ('Fulfillment', 'The next unlock — the largest remaining source of manual labor in the request lifecycle, and the direct path to an efficient ROI process', PRIMARY, PRIMARY_LT),
    ]
    for idx, (num, desc, num_color, bg) in enumerate(stat_data):
        if idx < len(b['stats']):
            sc = b['stats'][idx]
            add_rect(s, sc['x'], sc['y'], sc['w'], sc['h'],
                     fill=bg, line_color=BORDER, lw=0.5)
            add_text(s, num,  sc['x']+10, sc['y']+8,  sc['w']-16, 20,
                     size=14, bold=False, color=num_color)
            add_text(s, desc, sc['x']+10, sc['y']+30, sc['w']-16, sc['h']-36,
                     size=7, color=BODY, wrap=True)

    # Wordmark
    add_text(s, 'datavant', 900, 525, 50, 12, size=7.5, color=DIMMED, align=PP_ALIGN.RIGHT)
    return s


# ── SLIDE 3: The Problem ──────────────────────────────────────────────────────

async def build_slide3(prs, page, layout):
    print('  Slide 3 — The Problem …')
    path = BASE / 'slide_03_problem.html'
    b = await boxes(page, path, {
        'lbl':   '.slide-lbl',
        'h1':    '.h1',
        'sub':   '.subtitle',
        'stats': '.stat-card',
    })

    # Diagram region: just below subtitle → just above stats bar
    sub_bot   = b['sub'][0]['y'] + b['sub'][0]['h'] if b['sub'] else 70
    stats_top = b['stats'][0]['y'] if b['stats'] else 455

    out = SHOTS / 'slide_03_diag.png'
    diag_h = await clip_shot(page, path, sub_bot + 4, stats_top - 4, out)

    s = prs.slides.add_slide(layout)
    _slide_bg(s)

    # Header
    _header_lbl(s, b['lbl'])
    _header_h1(s, b['h1'])
    if b['sub']:
        add_rect(s, b['sub'][0]['x'], b['sub'][0]['y'], 3, b['sub'][0]['h'] + 4, fill=RED)
        add_text(s, b['sub'][0]['text'],
                 b['sub'][0]['x']+7, b['sub'][0]['y'],
                 b['sub'][0]['w'], b['sub'][0]['h']+8,
                 size=7.5, color=DIMMED, wrap=True)

    # Diagram image
    add_image(s, out, 0, sub_bot + 4, 960, diag_h)

    # Stat cards with left-color borders
    stat_data = [
        ('6.9 min',        'avg manual fulfillment — largest human time sink per request',                                                     TEAL,   rgb('#0d5c63')),
        ('≈ 0 min saved',  'partial fulfillment (7.07 min) saves no time vs manual (6.9 min)',                                                 AMBER,  AMBER_BD),
        ('23.5%',          'of digital requests fall into partial — zero efficiency gain',                                                     RED,    RED),
        ('54.5%',          'of partial are false partial — digital records retrieved, then discarded and redone manually',                      ORANGE, ORANGE),
    ]
    for idx, (num, desc, num_color, border_c) in enumerate(stat_data):
        if idx < len(b['stats']):
            sc = b['stats'][idx]
            # Left border accent
            add_rect(s, sc['x'], sc['y'], 3.5, sc['h'], fill=border_c)
            # Card body
            add_rect(s, sc['x']+3.5, sc['y'], sc['w']-3.5, sc['h'],
                     fill=SURF, line_color=BORDER, lw=0.5)
            # Text
            add_text(s, num,  sc['x']+9, sc['y']+5,  sc['w']-16, 18,
                     size=11.5, bold=False, color=num_color)
            add_text(s, desc, sc['x']+9, sc['y']+24, sc['w']-16, sc['h']-28,
                     size=6.5, color=BODY, wrap=True)

    # Sources line
    sources = ('Sources: 6.9 min — FTI Time Study Oct–Nov 2025 (n=2,271) · '
               '7.07 min partial, 1.83 min digital — Sigma DAR + HealthSource 2026 YTD · '
               '54.5% false partial — Digital ROI Performance Analysis Mar 2026 (n=165,267)')
    src_y = b['stats'][0]['y'] + b['stats'][0]['h'] + 4 if b['stats'] else 515
    add_text(s, sources, 24, src_y, 912, 16, size=5.5, color=DIMMED, wrap=True)
    add_text(s, 'datavant', 900, 525, 50, 12, size=7.5, color=DIMMED, align=PP_ALIGN.RIGHT)
    return s


# ── SLIDE 4: Proposed Hypothesis ─────────────────────────────────────────────

async def build_slide4(prs, page, layout):
    print('  Slide 4 — Proposed Hypothesis …')
    path = BASE / 'slide_04_vision.html'
    b = await boxes(page, path, {
        'lbl':      '.slide-lbl',
        'h1':       '.h1',
        'sub':      '.subtitle',
        'benefits': '.benefit-card',
    })

    sub_bot      = b['sub'][0]['y'] + b['sub'][0]['h'] if b['sub'] else 70
    benefits_top = b['benefits'][0]['y'] if b['benefits'] else 455

    out = SHOTS / 'slide_04_diag.png'
    diag_h = await clip_shot(page, path, sub_bot + 4, benefits_top - 4, out)

    s = prs.slides.add_slide(layout)
    _slide_bg(s)

    _header_lbl(s, b['lbl'])
    _header_h1(s, b['h1'])
    if b['sub']:
        add_rect(s, b['sub'][0]['x'], b['sub'][0]['y'], 3, b['sub'][0]['h'] + 4, fill=PRIMARY)
        add_text(s, b['sub'][0]['text'],
                 b['sub'][0]['x']+7, b['sub'][0]['y'],
                 b['sub'][0]['w'], b['sub'][0]['h']+8,
                 size=7.5, color=DIMMED, wrap=True)

    add_image(s, out, 0, sub_bot + 4, 960, diag_h)

    benefit_data = [
        ('Speed',      'Digital tasks automated, manual tasks routed simultaneously — no sequential waiting',                                              PRIMARY,    PRIMARY),
        ('Quality',    'Each worker handles one task type — no context switching, fewer errors',                                                            BODY,       BODY),
        ('Cost',       'Right worker, right task — standardized work is easier to train and lower cost',                                                   GRAY_MID,   GRAY_LIGHT),
        ('Visibility', 'Each task tracked separately — measure which task types are most labor-intensive to focus automation',                             SUCCESS,    SUCCESS),
    ]
    for idx, (title, desc, title_color, border_c) in enumerate(benefit_data):
        if idx < len(b['benefits']):
            bc = b['benefits'][idx]
            add_rect(s, bc['x'], bc['y'], 3.5, bc['h'], fill=border_c)
            add_rect(s, bc['x']+3.5, bc['y'], bc['w']-3.5, bc['h'],
                     fill=SURF, line_color=BORDER, lw=0.5)
            add_text(s, title, bc['x']+9, bc['y']+5, bc['w']-16, 15,
                     size=8, bold=True, color=title_color)
            add_text(s, desc,  bc['x']+9, bc['y']+20, bc['w']-16, bc['h']-24,
                     size=6.5, color=BODY, wrap=True)

    add_text(s, 'datavant', 900, 525, 50, 12, size=7.5, color=DIMMED, align=PP_ALIGN.RIGHT)
    return s


# ── SLIDE 5: Proposed Solution ────────────────────────────────────────────────

async def build_slide5(prs, page, layout):
    print('  Slide 5 — Proposed Solution …')
    path = BASE / 'slide_05_solution.html'
    b = await boxes(page, path, {
        'lbl':     '.slide-lbl',
        'h1':      '.h1',
        'flow':    '.flow-section',
        'build':   '.build-section',
        'blbl':    '.build-label',
        'bc':      '.bc',
        'bc5':     '.bc5',
        'bc5item': '.bc5-item',
    })

    # Screenshot the entire flow section
    if b['flow']:
        fl = b['flow'][0]
        out = SHOTS / 'slide_05_diag.png'
        await clip_shot(page, path, fl['y'], fl['y'] + fl['h'], out)
    else:
        out = None

    s = prs.slides.add_slide(layout)
    _slide_bg(s)

    # Header bar
    hdr_bg_h = 40
    if b['lbl'] and b['h1']:
        lbl_el = b['lbl'][0]
        h1_el  = b['h1'][0]
        # Light separator matches the .hdr border-bottom
        add_rect(s, 20, lbl_el['y'] - 2, 920, 38, fill=WHITE, line_color=BORDER, lw=0.5)
        add_text(s, lbl_el['text'], lbl_el['x'], lbl_el['y'], 120, lbl_el['h']+4,
                 size=7.5, color=PRIMARY)
        add_text(s, h1_el['text'], h1_el['x'], h1_el['y'], h1_el['w']+40, h1_el['h']+8,
                 size=12.75, bold=True, color=BODY)

    # Flow section image
    if out and b['flow']:
        fl = b['flow'][0]
        add_image(s, out, 0, fl['y'], 960, fl['h'])

    # Build section label
    if b['blbl']:
        bl = b['blbl'][0]
        add_text(s, 'What we need to build', bl['x'], bl['y'], 300, bl['h']+4,
                 size=5.5, bold=True, color=BODY)
        # Horizontal rule after label
        add_rect(s, bl['x'] + 160, bl['y'] + 4, 760, 1, fill=BORDER)

    # Build cards 1–4
    bc_data = [
        ('01', 'Record Location Registry',
         ['Health system map: systems, record types, date ranges',
          'Digital eligibility by type — powers task splitting'],
         TEAL, TEAL_LT, TEAL),
        ('02', 'Fulfillment Task Engine',
         ['Split by system × type × date range; route to right queue',
          'Leverages ROI Automation structured data'],
         DARK_NAVY, rgb('#f0f1f2'), DARK_NAVY),
        ('03', 'Task Assignment & Pickup UI',
         ['Worker interface to browse and claim assigned tasks',
          'Filter by type: onsite specialist vs. remote worker'],
         PRIMARY_800, rgb('#eef0fb'), PRIMARY_800),
        ('04', 'Medical Records Re-combination',
         ['Merge records from all tasks into a complete chart',
          'AOD assembly — per-task traceability for compliance'],
         PRIMARY, PRIMARY_LT, PRIMARY),
    ]
    for idx, (num, title, pts, accent, bg, txt_color) in enumerate(bc_data):
        if idx < len(b['bc']):
            bc = b['bc'][idx]
            # Top accent bar
            add_rect(s, bc['x'], bc['y'], bc['w'], 2.5, fill=accent)
            # Card body
            add_rect(s, bc['x'], bc['y']+2.5, bc['w'], bc['h']-2.5,
                     fill=bg, line_color=BORDER, lw=0.5)
            add_text(s, num,   bc['x']+6, bc['y']+5,   bc['w']-10, 10, size=5.5, bold=True, color=txt_color)
            add_text(s, title, bc['x']+6, bc['y']+14,  bc['w']-10, 18, size=6.75, bold=True, color=txt_color, wrap=True)
            y_pt = bc['y'] + 30
            for pt_txt in pts:
                add_text(s, '– ' + pt_txt, bc['x']+6, y_pt, bc['w']-10, 16, size=5.5, color=DIMMED, wrap=True)
                y_pt += 15

    # Card 5 — full width
    if b['bc5']:
        c5 = b['bc5'][0]
        add_rect(s, c5['x'], c5['y'], c5['w'], 2.5, fill=GRAY_MID)
        add_rect(s, c5['x'], c5['y']+2.5, c5['w'], c5['h']-2.5,
                 fill=rgb('#f8fafc'), line_color=rgb('#9ca3af'), lw=0.5)
        add_text(s, '05', c5['x']+8, c5['y']+5, 30, 10, size=5.5, bold=True, color=GRAY_MID)
        add_text(s, 'Tracking, Instrumentation & Resilience',
                 c5['x']+28, c5['y']+5, 400, 12, size=6.75, bold=True, color=rgb('#374151'))

        # 3 sub-panels
        sub_items = [
            ('Audit Trails',          'Per-task record: system, worker, what retrieved, when — compliance and dispute resolution',            GRAY_LIGHT),
            ('Data Instrumentation',  'Time-on-task by type, system, and worker — enables true labor cost measurement and automation targeting', PRIMARY),
            ('Exception Handling',    'Graceful fallback when digital fails — reroutes that task to manual without blocking other tasks',      AMBER_BD),
        ]
        body_x = c5['x'] + 6
        body_y = c5['y'] + 20
        body_h = c5['h'] - 24
        panel_w = (c5['w'] - 18) / 3
        for pi, (ptitle, pdesc, pleft) in enumerate(sub_items):
            px = body_x + pi * (panel_w + 5)
            add_rect(s, px, body_y, 2.5, body_h, fill=pleft)
            add_rect(s, px+2.5, body_y, panel_w-2.5, body_h, fill=SURF)
            add_text(s, ptitle, px+7, body_y+4, panel_w-12, 12, size=6.5, bold=True, color=BODY)
            add_text(s, pdesc,  px+7, body_y+16, panel_w-12, body_h-20, size=5.5, color=DIMMED, wrap=True)

    add_text(s, 'datavant', 900, 525, 50, 12, size=7.5, color=DIMMED, align=PP_ALIGN.RIGHT)
    return s


# ── SLIDE 6: Roadmap ─────────────────────────────────────────────────────────

async def build_slide6(prs, page, layout):
    print('  Slide 6 — Roadmap …')
    path = BASE / 'slide_06_roadmap.html'
    b = await boxes(page, path, {
        'lbl':     '.slide-lbl',
        'h1':      '.h1',
        'sub':     '.subtitle',
        'roadmap': '.roadmap',
    })

    # Screenshot the entire roadmap phase area
    if b['roadmap']:
        rm = b['roadmap'][0]
        out = SHOTS / 'slide_06_diag.png'
        # Add a little padding above
        await clip_shot(page, path, rm['y'] - 2, rm['y'] + rm['h'], out)
        diag_y = rm['y'] - 2
        diag_h = rm['h'] + 2
    else:
        out = None

    s = prs.slides.add_slide(layout)
    _slide_bg(s)

    _header_lbl(s, b['lbl'])
    _header_h1(s, b['h1'])
    if b['sub']:
        sub = b['sub'][0]
        add_rect(s, sub['x'], sub['y'], 3, sub['h'] + 4, fill=PRIMARY)
        add_text(s, sub['text'], sub['x']+7, sub['y'], sub['w'], sub['h']+8,
                 size=7.5, color=DIMMED, wrap=True)

    if out and b['roadmap']:
        add_image(s, out, 0, diag_y, 960, diag_h)

    add_text(s, 'datavant', 900, 525, 50, 12, size=7.5, color=DIMMED, align=PP_ALIGN.RIGHT)
    return s


# ── Main ──────────────────────────────────────────────────────────────────────

async def build_slide_exec(prs, page, layout):
    print('  Slide 2 — Executive Summary (full image) …')
    path = BASE / 'slide_02_exec_summary.html'
    out  = SHOTS / 'slide_02_exec.png'
    await page.goto(f'file://{path}')
    await page.wait_for_timeout(1200)
    await page.screenshot(path=str(out),
                          clip={'x': 0, 'y': TOP, 'width': 960, 'height': 540})
    s = prs.slides.add_slide(layout)
    add_image(s, out, 0, 0, 960, 540)
    return s


async def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(
            viewport={'width': 960, 'height': 540 + TOP + 40}
        )

        await build_slide1(prs, page, blank)
        await build_slide_exec(prs, page, blank)
        await build_slide2(prs, page, blank)
        await build_slide3(prs, page, blank)
        await build_slide4(prs, page, blank)
        await build_slide5(prs, page, blank)
        await build_slide6(prs, page, blank)

        await browser.close()

    prs.save(str(OUT))
    print(f'\nSaved → {OUT}')
    print(f'Slides: {len(prs.slides)}  |  All flow diagrams preserved as images.')
    print('Text elements (labels, titles, stat cards, build cards) are fully editable.')

asyncio.run(main())
