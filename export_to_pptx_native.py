#!/usr/bin/env python3
"""
Fully native PPTX export — all elements editable. Reflects HTML slides June 18, 2026.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUT = Path('/Users/jennywang/Documents/Claude/New Fulfillment Workflow/New Fulfillment Flow Vision.pptx')
SCALE = 10.0 / 960

def i(px): return Inches(px * SCALE)

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:6],16))

PRIMARY     = rgb('#2945f0')
PRIMARY_800 = rgb('#102cd5')
PRIMARY_LT  = rgb('#eceefe')
BODY        = rgb('#020202')
DIMMED      = rgb('#687382')
BORDER      = rgb('#c2c7ce')
SURF        = rgb('#f2f3f8')
SURF_FA     = rgb('#f8f9fa')
WHITE       = rgb('#ffffff')
BLACK       = rgb('#000000')
RED         = rgb('#cc3300')
RED_LT      = rgb('#fff1ef')
AMBER       = rgb('#b45309')
AMBER_BD    = rgb('#f59e0b')
AMBER_LT    = rgb('#fef3c7')
GREEN       = rgb('#1a7f4b')
DARK_NAVY   = rgb('#1a1a2e')
GRAY_MID    = rgb('#6b7280')
GRAY_LIGHT  = rgb('#9ca3af')
CERNER_BG   = rgb('#1e3a5f')
EPIC_BG     = rgb('#1a1a2e')
BATCH_BG    = rgb('#374151')
IMAGING_BG  = rgb('#7c2d12')
ARCHIVAL_BG = rgb('#44403c')
TEAL        = rgb('#1a7f4b')

def _no_line(shp):  shp.line.fill.background()
def _solid(shp, c): shp.fill.solid(); shp.fill.fore_color.rgb = c
def _no_fill(shp):  shp.fill.background()

def _no_inset(tf):
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        for a in ('lIns','tIns','rIns','bIns'):
            bp.set(a, '0')

def add_rect(slide, x, y, w, h, fill=None, lc=None, lw=0.75):
    shp = slide.shapes.add_shape(1, i(x), i(y), i(w), i(h))
    if fill: _solid(shp, fill)
    else:    _no_fill(shp)
    if lc:
        shp.line.color.rgb = lc
        shp.line.width = Pt(lw)
    else:
        _no_line(shp)
    return shp

def add_dashed(slide, x, y, w, h, lc, fill=None, lw=1.5):
    shp = add_rect(slide, x, y, w, h, fill=fill, lc=lc, lw=lw)
    ln = shp._element.spPr.find(qn('a:ln'))
    if ln is not None:
        for pd in ln.findall(qn('a:prstDash')): ln.remove(pd)
        pd = etree.SubElement(ln, qn('a:prstDash'))
        pd.set('val', 'dash')
    return shp

def add_text(slide, text, x, y, w, h,
             sz=9, bold=False, italic=False, color=BODY,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    for idx, line in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name   = 'Geist'
        run.font.size   = Pt(sz)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
    _no_inset(tf)
    return txb

def wm(slide):
    add_text(slide, 'datavant', 896, 524, 56, 12, sz=7.5, color=DIMMED, align=PP_ALIGN.RIGHT)


# ── SLIDE 1: Cover ────────────────────────────────────────────────────────────

def build_slide1(prs, layout):
    print('  Slide 1 — Cover')
    s = prs.slides.add_slide(layout)
    add_rect(s, 0, 0, 960, 540, fill=BLACK)
    add_rect(s, 560, 60, 420, 420, fill=rgb('#091877'))  # orb
    add_text(s, 'datavant', 48, 36, 140, 20, sz=13, bold=True, color=WHITE)
    add_text(s, 'New Fulfillment Workflow  ·  Proposed Hypothesis',
             48, 196, 560, 14, sz=8.5, color=DIMMED)
    add_text(s, 'Efficient ROI', 48, 216, 560, 44, sz=38, bold=True, color=WHITE)
    add_text(s, 'Fulfillment', 48, 258, 560, 40, sz=35, color=WHITE)
    add_text(s, 'Replacing sequential, context-switching fulfillment with automated parallel task routing — digital and manual components completed simultaneously.',
             48, 312, 500, 52, sz=11, color=rgb('#7f8f9f'), wrap=True)
    add_text(s, 'June 2026', 48, 490, 140, 14, sz=9, color=DIMMED)
    return s


# ── SLIDE 2: ROI North Star ───────────────────────────────────────────────────

def build_slide2(prs, layout):
    print('  Slide 2 — ROI North Star')
    s = prs.slides.add_slide(layout)
    add_rect(s, 0, 0, 960, 540, fill=WHITE)
    P, H = 36, 28

    add_text(s, 'ROI North Star', P, H, 200, 12, sz=7.5, color=PRIMARY)
    add_text(s, 'An Efficient, Cost-Effective ROI Process — Rooted in Automation',
             P, H+14, 888, 26, sz=17, color=BODY)
    add_text(s, 'Every phase reduces redundant labor and drives cost out. Phases 1–3 addressed logging and auth. H2 2026 completes the picture — efficient fulfillment and automated QC.',
             P, H+42, 888, 20, sz=9, color=DIMMED, wrap=True)

    py = H + 66
    add_text(s, 'The Request Lifecycle — Automating toward efficiency at every step',
             P, py, 600, 12, sz=7, color=DIMMED)

    arr_w = 18
    bar_y = py + 14
    bh = 44
    bw = (888 - 5*arr_w) // 6
    steps = [
        ('INTAKE',       SURF,      DIMMED,  False),
        ('LOGGING',      PRIMARY_LT, PRIMARY, False),
        ('AUTH\nREVIEW', PRIMARY_LT, PRIMARY, False),
        ('FULFILLMENT',  PRIMARY,    WHITE,   True),
        ('QC',           PRIMARY,    WHITE,   True),
        ('DELIVERY',     SURF,      DIMMED,  False),
    ]
    for idx, (label, bg, tc, is_bold) in enumerate(steps):
        cx = P + idx * (bw + arr_w)
        add_rect(s, cx, bar_y, bw, bh, fill=bg, lc=BORDER, lw=0.5)
        add_text(s, label, cx+2, bar_y+4, bw-4, bh-8,
                 sz=8.5, bold=is_bold, color=tc, align=PP_ALIGN.CENTER, wrap=True)
        if idx < 5:
            add_text(s, '›', cx+bw, bar_y+14, arr_w, 16, sz=9, color=DIMMED, align=PP_ALIGN.CENTER)

    # Badge row
    badge_y = bar_y + bh + 6
    inflight_x = P + bw + arr_w
    inflight_w = bw + arr_w + bw
    add_rect(s, inflight_x, badge_y, inflight_w, 20, fill=PRIMARY_LT, lc=PRIMARY, lw=0.75)
    add_text(s, 'In Flight — ROI Automation Phases 1–3 · Piloting July 2026',
             inflight_x+4, badge_y+3, inflight_w-8, 14, sz=7.5, color=PRIMARY, align=PP_ALIGN.CENTER)

    fulf_x = P + 3*(bw + arr_w)
    add_rect(s, fulf_x, badge_y, bw, 20, fill=PRIMARY)
    add_text(s, 'H2 2026', fulf_x+2, badge_y+2, bw-4, 8,
             sz=6, bold=True, color=rgb('#ffffffb3'), align=PP_ALIGN.CENTER)
    add_text(s, 'New Fulfillment Workflow', fulf_x+2, badge_y+10, bw-4, 9,
             sz=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    qc_x = P + 4*(bw + arr_w)
    add_rect(s, qc_x, badge_y, bw, 20, fill=PRIMARY)
    add_text(s, 'H2 2026', qc_x+2, badge_y+2, bw-4, 8,
             sz=6, bold=True, color=rgb('#ffffffb3'), align=PP_ALIGN.CENTER)
    add_text(s, 'QC Automation', qc_x+2, badge_y+10, bw-4, 9,
             sz=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    div_y = badge_y + 28
    add_rect(s, P, div_y, 888, 1, fill=BORDER)
    add_text(s, 'Phases 1–3 lay the foundation — accurate, structured data at every request. H2 2026 targets the largest remaining manual stage: fulfillment and QC together.',
             P, div_y+6, 888, 24, sz=9, italic=True, color=DIMMED, wrap=True)

    stat_y = div_y + 38
    stat_h = 540 - stat_y - 22
    stat_w = (888 - 16) // 3
    stats = [
        ('2 min → 50 sec', PRIMARY,
         'processing time per request from Phases 1–3 — AI extraction reducing manual logging and auth work at scale',
         SURF, BORDER),
        ('32M requests', PRIMARY,
         'processed in 2025 — every minute saved across fulfillment and QC at this volume = significant labor cost out',
         SURF, BORDER),
        ('Fulfillment', PRIMARY_800,
         'The largest remaining source of manual labor in the request lifecycle — and the direct path to an efficient, cost-effective ROI process',
         PRIMARY_LT, PRIMARY),
    ]
    for idx, (num, nc, desc, bg, bc) in enumerate(stats):
        sx = P + idx * (stat_w + 8)
        add_rect(s, sx, stat_y, stat_w, stat_h, fill=bg, lc=bc, lw=0.75)
        add_text(s, num, sx+10, stat_y+8, stat_w-16, 22, sz=17, color=nc)
        if idx == 2:
            add_text(s, 'THE NEXT UNLOCK', sx+10, stat_y+30, stat_w-16, 10,
                     sz=6.5, bold=True, color=PRIMARY)
        add_text(s, desc, sx+10, stat_y+(32 if idx==2 else 28),
                 stat_w-16, stat_h-36, sz=7.5, color=BODY, wrap=True)

    wm(s)
    return s


# ── SLIDE 3: The Problem ──────────────────────────────────────────────────────

def build_slide3(prs, layout):
    print('  Slide 3 — The Problem')
    s = prs.slides.add_slide(layout)
    add_rect(s, 0, 0, 960, 540, fill=WHITE)
    P, H = 24, 18

    add_text(s, 'The Problem', P, H, 200, 12, sz=7.5, color=RED)
    add_text(s, 'One Request. Five Systems. One Person Doing It All — Sequentially.',
             P, H+14, 912, 22, sz=13.5, bold=True, color=BODY)
    add_rect(s, P, H+38, 3, 22, fill=RED)
    add_text(s, 'Today, a single operations staff member must be trained on 5 different systems and context-switch between all of them to fulfill one request. Every login, every walk across the hospital, every phone call to retrieve records — all on one person, one task at a time.',
             P+7, H+38, 900, 22, sz=7.5, color=DIMMED, wrap=True)

    # Example strip
    ex_y = H + 64
    ex_h = 48
    add_rect(s, P, ex_y, 912, ex_h, fill=SURF, lc=BORDER, lw=0.75)
    add_text(s, 'Example request', P+10, ex_y+6, 90, 10, sz=7.5, bold=True, color=DIMMED)
    add_text(s, 'Bomi Kim, 01/01/1990\nItemized Bill, MRI CD, Radiology, Progress notes, Mammography (1990–present)',
             P+10, ex_y+17, 190, 28, sz=7.5, color=BODY, wrap=True)
    add_rect(s, P+208, ex_y+9, 1, 30, fill=BORDER)
    add_text(s, '→', P+212, ex_y+16, 16, 16, sz=12, color=DIMMED, align=PP_ALIGN.CENTER)
    add_text(s, 'Splits across 5 systems:', P+232, ex_y+6, 120, 10, sz=7.5, bold=True, color=DIMMED)

    sx = P + 366
    systems_ex = [
        ('Epic 2025–now',       [('↓ Radiology', True), ('✋ Mammo', False)]),
        ('Cerner 2021–25',      [('↓ Radiology', True), ('✋ Mammo', False), ('✋ Cardio', False)]),
        ('Imaging System',            [('✋ MRI CD', False)]),
        ('Archival System',           [('✋ 1990–2021', False)]),
        ('PM System',                 [('✋ Bill', False)]),
    ]
    for (name, tasks) in systems_ex:
        sw = 88
        add_rect(s, sx, ex_y+5, sw, ex_h-10, fill=WHITE, lc=BORDER, lw=0.5)
        add_text(s, name, sx+4, ex_y+7, sw-8, 11, sz=6, bold=True, color=BODY)
        ty = ex_y+20
        for (task, is_dig) in tasks:
            tc = rgb('#1e40af') if is_dig else rgb('#92400e')
            tf = rgb('#dbeafe') if is_dig else rgb('#fef9c3')
            tb = rgb('#93c5fd') if is_dig else rgb('#fde047')
            add_rect(s, sx+3, ty, sw-8, 9, fill=tf, lc=tb, lw=0.3)
            add_text(s, task, sx+5, ty+1, sw-12, 8, sz=5.5, color=tc)
            ty += 10
        sx += sw + 5

    # Flow row — compute height to leave fixed space for pain bar
    pb_h = 52
    pb_y = 540 - pb_h - 14
    fl_y = ex_y + ex_h + 8
    fl_h = pb_y - 8 - fl_y

    # Worker badge
    wb_w = 54
    add_rect(s, P, fl_y, wb_w, fl_h, fill=SURF, lc=BORDER, lw=0.5)
    add_text(s, '■■', P+wb_w//2-8, fl_y+10, 20, 14, sz=10, color=DIMMED, align=PP_ALIGN.CENTER)
    add_text(s, 'Ops\nStaff', P+3, fl_y+28, wb_w-6, 18, sz=7.5, bold=True, color=BODY, align=PP_ALIGN.CENTER, wrap=True)
    add_text(s, 'trained\non all 5', P+3, fl_y+48, wb_w-6, 16, sz=7, color=RED, align=PP_ALIGN.CENTER, wrap=True)

    arr_entry = 22
    add_text(s, '→', P+wb_w, fl_y + fl_h//2 - 8, arr_entry, 16, sz=12, color=BORDER, align=PP_ALIGN.CENTER)

    # 5 system cards
    ctx_w = 32
    repeat_w = 60
    arr_end = 22
    avail = 912 - wb_w - arr_entry - 5*ctx_w - repeat_w - arr_end
    card_w = avail // 5

    card_defs = [
        ('Cerner',          '2021–2025',   'Digital',   CERNER_BG,
         ['Login to Cerner', 'Retrieve records', 'Upload to portal'], 'Multiple steps', False),
        ('Epic',            '2025–now',    'Digital',   EPIC_BG,
         ['Login to Epic',   'Retrieve records', 'Upload to portal'], 'Multiple steps', False),
        ('PM System',       'Billing records',  'Digital',   BATCH_BG,
         ['Login to PM System', 'Retrieve itemized bill', 'Upload to portal'], 'Multiple steps', False),
        ('Imaging\nSystem', 'Radiology dept',   'Physical',  IMAGING_BG,
         ['Walk across hospital', 'Wait for department', 'Walk back to desk'], 'Up to 30 min', True),
        ('Archival\nSystem','Records 1990–2021', 'Off-site', ARCHIVAL_BG,
         ['Call to request boxes', 'Wait for delivery', 'Search records'], 'Hours to days', True),
    ]
    ctx_labels = ['logout &\nswitch', 'logout &\nswitch', 'leave desk\n& walk', 'call out\n& wait']
    ctx_colors = [RED, RED, AMBER, AMBER]

    cx = P + wb_w + arr_entry
    for cidx, (sname, sdate, stype, sbg, steps_list, time_txt, time_warn) in enumerate(card_defs):
        add_rect(s, cx, fl_y, card_w, fl_h, fill=sbg)
        add_text(s, sname, cx+6, fl_y+8, card_w-12, 22, sz=9, bold=True, color=WHITE, wrap=True)
        add_text(s, sdate, cx+6, fl_y+30, card_w-12, 10, sz=6.5, color=rgb('#ffffff66'), wrap=True)

        type_colors = {
            'Digital':  (rgb('#3b5bfb'), rgb('#ffffff99')),
            'Physical': (rgb('#dc2626'), rgb('#fca5a5')),
            'Off-site': (rgb('#b45309'), rgb('#fcd34d')),
        }
        tbg, ttc = type_colors[stype]
        bw2 = max(38, len(stype)*6+8)
        add_rect(s, cx+6, fl_y+42, bw2, 12, fill=tbg)
        add_text(s, stype, cx+7, fl_y+43, bw2-2, 10, sz=6, bold=True, color=ttc)

        step_bar_fills = [
            rgb('#2d4f82'), rgb('#2d2d50'), rgb('#4a5570'),
            rgb('#8f3d22'), rgb('#5a5552'),
        ]
        sy = fl_y + 58
        for step in steps_list:
            add_rect(s, cx+4, sy, card_w-8, 20, fill=step_bar_fills[cidx])
            add_text(s, step, cx+7, sy+3, card_w-14, 15, sz=6.5, color=WHITE, wrap=True)
            sy += 23

        time_c = rgb('#fca5a5') if time_warn else rgb('#ffffff66')
        add_text(s, time_txt, cx+6, fl_y + fl_h - 18, card_w-12, 14, sz=9, bold=True, color=time_c)

        if cidx < 4:
            ctx_c = ctx_colors[cidx]
            ctx_x = cx + card_w
            add_text(s, '⇄', ctx_x+2, fl_y + fl_h//2 - 12, ctx_w-4, 14,
                     sz=9, bold=True, color=ctx_c, align=PP_ALIGN.CENTER)
            add_text(s, ctx_labels[cidx], ctx_x, fl_y + fl_h//2 + 4, ctx_w, 18,
                     sz=5.5, bold=True, color=ctx_c, align=PP_ALIGN.CENTER, wrap=True)
            cx += card_w + ctx_w
        else:
            cx += card_w

    add_text(s, '→', cx, fl_y + fl_h//2 - 8, arr_end, 16, sz=11, color=BORDER, align=PP_ALIGN.CENTER)
    rep_x = cx + arr_end
    add_dashed(s, rep_x, fl_y + fl_h//4, repeat_w, fl_h//2, lc=RED, fill=WHITE, lw=1.5)
    add_text(s, 'Repeat for\nevery request', rep_x+4, fl_y + fl_h//4 + 4, repeat_w-8, 20,
             sz=6.5, bold=True, color=RED, align=PP_ALIGN.CENTER, wrap=True)
    add_text(s, 'same steps,\nevery time', rep_x+4, fl_y + fl_h//4 + 26, repeat_w-8, 18,
             sz=6, color=DIMMED, align=PP_ALIGN.CENTER, wrap=True)

    # Pain bar
    pain_w = (912 - 3*5) // 4
    pains = [
        ('5 System Logins',
         'Operations staff must be trained on and log into every system — credentials, UIs, and workflows for all 5'),
        ('Constant Context Switching',
         'Every system switch breaks focus — operations staff must mentally reset and reorient for each new environment'),
        ('Up to 30 Min Physical Walk',
         'Imaging records require leaving the desk, walking across the hospital, waiting, and walking back — every single time'),
        ('Archival Delays Hours–Days',
         'Records from 1990 require calling out, waiting for physical boxes to arrive, then manually searching through them'),
    ]
    for idx, (lbl, desc) in enumerate(pains):
        px2 = P + idx * (pain_w + 5)
        add_rect(s, px2, pb_y, 3, pb_h, fill=RED)
        add_rect(s, px2+3, pb_y, pain_w-3, pb_h, fill=RED_LT, lc=rgb('#fecdc8'), lw=0.5)
        add_text(s, lbl,  px2+8, pb_y+5,  pain_w-14, 12, sz=8.5, bold=True, color=RED)
        add_text(s, desc, px2+8, pb_y+18, pain_w-14, pb_h-22, sz=7, color=rgb('#7f2f20'), wrap=True)

    wm(s)
    return s


# ── SLIDE 4: Proposed Hypothesis ─────────────────────────────────────────────

def build_slide4(prs, layout):
    print('  Slide 4 — Proposed Hypothesis')
    s = prs.slides.add_slide(layout)
    add_rect(s, 0, 0, 960, 540, fill=WHITE)
    P, H = 24, 18

    add_text(s, 'Proposed Hypothesis', P, H, 200, 12, sz=7.5, color=PRIMARY)
    add_text(s, 'One Engine. Specialized Processors. No Context Switching.',
             P, H+14, 912, 22, sz=13.5, bold=True, color=BODY)
    add_rect(s, P, H+38, 3, 28, fill=PRIMARY)
    add_text(s, "The fulfillment engine aggregates tasks by type across all patients — so each processor works only what they’re trained on, stays in one system all day, and completes work for 10 patients at once instead of one. No login/logout. No walking to a different location. No switching between patients mid-task.",
             P+7, H+38, 900, 28, sz=7.5, color=DIMMED, wrap=True)

    # Before strip
    bs_y = H + 70
    bs_h = 30
    add_rect(s, P, bs_y, 912, bs_h, fill=SURF_FA, lc=BORDER, lw=0.75)
    add_text(s, 'TODAY', P+10, bs_y+8, 46, 12, sz=7.5, bold=True, color=GRAY_LIGHT)
    before_steps = ['Login Epic (Pt 1)', 'Switch to Cerner', 'Request from Archival',
                    'Login Epic (Pt 2), repeat', 'Cerner (Pt 2)', '…']
    bx = P + 64
    for step in before_steps:
        sw = max(48, len(step)*5+10)
        add_rect(s, bx, bs_y+7, sw, 16, fill=rgb('#f3f4f6'), lc=BORDER, lw=0.5)
        add_text(s, step, bx+3, bs_y+9, sw-6, 12, sz=6.5, color=GRAY_LIGHT)
        bx += sw + 14
        if bx < P + 760:
            add_text(s, '→', bx-12, bs_y+9, 10, 12, sz=8, color=rgb('#d1d5db'), align=PP_ALIGN.CENTER)
    add_rect(s, 866, bs_y+7, 70, 16, fill=rgb('#fef2f2'), lc=rgb('#fecaca'), lw=0.5)
    add_text(s, '⚠ Context switching kills throughput', 868, bs_y+9, 66, 12,
             sz=6, bold=True, color=rgb('#dc2626'), wrap=True)

    # Flow row
    ben_h = 46
    ben_y = 540 - ben_h - 14
    fl_y = bs_y + bs_h + 8
    fl_h = ben_y - 8 - fl_y

    # Patients column
    pt_w = 108
    add_text(s, 'REQUESTS', P, fl_y, pt_w, 10, sz=7, bold=True, color=DIMMED)
    pt_card_h = (fl_h - 14 - 4*3) // 5
    patients = [
        ('Patient 1', 'Epic · Cerner · Archival'),
        ('Patient 2', 'Epic · Cerner'),
        ('Patient 3', 'Cerner · Imaging'),
        ('Patient 4', 'Epic · Archival'),
        ('+ 6 more',  'various systems'),
    ]
    for pidx, (pname, psys) in enumerate(patients):
        py2 = fl_y + 14 + pidx * (pt_card_h + 3)
        bg = SURF if pidx < 4 else WHITE
        nc = BODY if pidx < 4 else DIMMED
        add_rect(s, P, py2, pt_w, pt_card_h, fill=bg, lc=BORDER, lw=0.5)
        add_text(s, pname, P+5, py2+3, pt_w-10, 12, sz=8, bold=True, color=nc)
        add_text(s, psys,  P+5, py2+14, pt_w-10, 12, sz=7, color=DIMMED)

    # Arrow + Engine
    arr_w2 = 18
    eng_x = P + pt_w + arr_w2
    eng_w = 118
    add_text(s, '→', P+pt_w, fl_y + fl_h//2 - 8, arr_w2, 16, sz=12, color=BORDER, align=PP_ALIGN.CENTER)
    add_rect(s, eng_x, fl_y, eng_w, fl_h, fill=BODY)
    add_text(s, 'Fulfillment\nEngine', eng_x+6, fl_y+12, eng_w-12, 30,
             sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
    add_text(s, 'Breaks down & aggregates\ntasks by system + type',
             eng_x+6, fl_y+46, eng_w-12, 36, sz=7, color=rgb('#ffffff8a'), align=PP_ALIGN.CENTER, wrap=True)
    add_rect(s, eng_x+10, fl_y+86, eng_w-20, 16, fill=rgb('#1a7f4b1f'), lc=rgb('#22c55e40'), lw=0.5)
    add_text(s, 'Intelligent routing', eng_x+12, fl_y+88, eng_w-24, 12,
             sz=6.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, '→', eng_x+eng_w, fl_y + fl_h//2 - 8, arr_w2, 16, sz=12, color=BORDER, align=PP_ALIGN.CENTER)

    # Processors
    proc_x = eng_x + eng_w + arr_w2
    proc_w = P + 912 - proc_x
    gap5 = 5
    proc_h = (fl_h - 2*gap5) // 3

    proc_defs = [
        ('Processor A', 'Cerner-trained only\nSingle login, all day', CERNER_BG,
         ['Cerner · Pt 1', 'Cerner · Pt 2', 'Cerner · Pt 3',
          'Cerner · Pt 6', 'Cerner · Pt 7', 'Cerner · Pt 9', '+ 3 more'],
         'Stays in Cerner. Completes 9 patients without switching systems once.'),
        ('Processor B', 'Epic-trained only\nSingle login, all day', EPIC_BG,
         ['Epic · Pt 1', 'Epic · Pt 2', 'Epic · Pt 4',
          'Epic · Pt 5', 'Epic · Pt 8', '+ 2 more'],
         'Stays in Epic. Handles 7 patients — no context switching, deep expertise.'),
        ('Processor C', 'Archival + Imaging\nAuto-aggregated batch', BATCH_BG,
         ['Archival System batch · 10 patients',
          'Imaging · Pt 3', 'Imaging · Pt 5', 'Imaging · Pt 8'],
         'One trip. Grabs archival records for 10 patients at once.'),
    ]
    outcome_w = 165
    for pidx, (pname, pspec, pbg, chips, outcome) in enumerate(proc_defs):
        py2 = fl_y + pidx * (proc_h + gap5)
        add_rect(s, proc_x, py2, proc_w, proc_h, fill=pbg)
        add_text(s, pname, proc_x+8, py2+6, 120, 14, sz=10, bold=True, color=WHITE)
        add_text(s, pspec, proc_x+8, py2+20, 120, 22, sz=7, color=rgb('#ffffff7f'), wrap=True, italic=True)
        chip_x = proc_x + 140
        chip_y = py2 + 8
        for chip in chips:
            is_big = 'batch' in chip.lower()
            cw = max(42, len(chip)*5+14)
            add_rect(s, chip_x, chip_y, cw, 14, fill=rgb('#5a6070' if is_big else '#3a4560'),
                     lc=rgb('#6a7080'), lw=0.5)
            add_text(s, chip, chip_x+4, chip_y+2, cw-8, 10, sz=6, bold=is_big, color=WHITE)
            chip_x += cw + 5
            if chip_x > proc_x + proc_w - outcome_w - 20:
                chip_x = proc_x + 140
                chip_y += 17
        add_rect(s, proc_x + proc_w - outcome_w - 5, py2+4, 1, proc_h-8,
                 fill=rgb('#ffffff30'))
        add_text(s, outcome, proc_x + proc_w - outcome_w, py2+6,
                 outcome_w-8, proc_h-12, sz=7.5, color=rgb('#ffffffa6'), wrap=True)

    # Benefits bar
    ben_w = (912 - 3*6) // 4
    benefits = [
        ('Stay in Flow',       PRIMARY,  PRIMARY,  'Each processor works one system, one login — no interruptions or location changes'),
        ('Batch Efficiency',   GREEN,    GREEN,    'Archival & Imaging tasks aggregate automatically — one trip handles 10 patients at once'),
        ('Deep Specialization',BODY,     BODY,     'Train each processor on one system — faster throughput, fewer errors, easier QA'),
        ('True Parallelism',   GRAY_MID, GRAY_MID, 'All processors run simultaneously — total time = slowest lane, not the sum of all'),
    ]
    for idx, (title, tc, bc, desc) in enumerate(benefits):
        bx = P + idx * (ben_w + 6)
        add_rect(s, bx, ben_y, 3, ben_h, fill=bc)
        add_rect(s, bx+3, ben_y, ben_w-3, ben_h, fill=SURF, lc=BORDER, lw=0.5)
        add_text(s, title, bx+8, ben_y+5, ben_w-14, 13, sz=9, bold=True, color=tc)
        add_text(s, desc,  bx+8, ben_y+18, ben_w-14, ben_h-22, sz=7, color=BODY, wrap=True)

    wm(s)
    return s


# ── SLIDE 5: Proposed Solution ────────────────────────────────────────────────

def build_slide5(prs, layout):
    print('  Slide 5 — Proposed Solution')
    s = prs.slides.add_slide(layout)
    add_rect(s, 0, 0, 960, 540, fill=WHITE)
    P, H = 24, 22

    add_text(s, 'Proposed Solution', P, H, 200, 12, sz=7.5, color=PRIMARY)
    add_text(s, 'How the Fulfillment Engine Works',
             P, H+14, 912, 24, sz=16.5, bold=True, color=BODY)
    add_rect(s, P, H+40, 3, 22, fill=PRIMARY)
    add_text(s, 'The new workflow slots into the existing ROI pipeline between intake and QC. The Task Engine reads from a Record Location Registry to intelligently split each request and route tasks in parallel.',
             P+7, H+40, 900, 22, sz=8, color=DIMMED, wrap=True)

    fl_y = H + 68
    fl_h = 240

    # Geometry
    intake_w = 112
    farr_w   = 24
    bound_w  = 540
    qc_w     = 88
    del_w    = 88
    intake_x = P
    farr1_x  = intake_x + intake_w
    bound_x  = farr1_x + farr_w
    farr2_x  = bound_x + bound_w
    qc_x     = farr2_x + farr_w
    farr3_x  = qc_x + qc_w
    del_x    = farr3_x + farr_w

    # Intake
    add_rect(s, intake_x, fl_y, intake_w, fl_h, fill=SURF, lc=BORDER, lw=0.75)
    add_text(s, 'REQUEST INTAKE', intake_x+6, fl_y+8, intake_w-12, 10,
             sz=6, bold=True, color=DIMMED)
    add_text(s, 'Turbo', intake_x+6, fl_y+20, intake_w-12, 16,
             sz=12, bold=True, color=BODY)
    add_text(s, 'Extracts fields, sorts to Site ID',
             intake_x+6, fl_y+38, intake_w-12, 28, sz=7.5, color=DIMMED, wrap=True)

    add_text(s, '→', farr1_x, fl_y + fl_h//2 - 8, farr_w, 16,
             sz=12, color=DIMMED, align=PP_ALIGN.CENTER)

    # Boundary (dashed)
    add_dashed(s, bound_x, fl_y, bound_w, fl_h, lc=rgb('#93c5fd'), fill=rgb('#f8f9ff'), lw=1.5)
    lbl_bx = bound_x + bound_w//2 - 80
    add_rect(s, lbl_bx, fl_y + fl_h - 13, 160, 14, fill=PRIMARY)
    add_text(s, 'New Fulfillment Workflow', lbl_bx+2, fl_y + fl_h - 12, 156, 12,
             sz=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Task Engine
    bp = 10
    eng_w2 = 148
    eng_x2 = bound_x + bp
    eng_inner_h = fl_h - 2*bp - 16
    add_rect(s, eng_x2, fl_y+bp, eng_w2, eng_inner_h, fill=DARK_NAVY)
    add_text(s, 'FULFILLMENT', eng_x2+6, fl_y+bp+6, eng_w2-12, 10,
             sz=5.5, bold=True, color=rgb('#ffffff66'))
    add_text(s, 'Task\nEngine', eng_x2+6, fl_y+bp+18, eng_w2-12, 28,
             sz=11, bold=True, color=WHITE, wrap=True)
    add_text(s, 'Splits request into task units. Routes each to digital or manual queue simultaneously.',
             eng_x2+6, fl_y+bp+48, eng_w2-12, 56, sz=7, color=rgb('#ffffff8a'), wrap=True)
    add_text(s, '↑ reads from Record Location Registry',
             eng_x2+6, fl_y+bp+108, eng_w2-12, 18, sz=6, bold=True, color=GREEN, wrap=True)

    ifr_x = eng_x2 + eng_w2 + 4
    add_text(s, '→', ifr_x, fl_y + fl_h//2 - 8, 16, 16, sz=10,
             color=rgb('#ffffff40'), align=PP_ALIGN.CENTER)

    # Tasks stacked
    tasks_x = ifr_x + 16
    tasks_w = 192
    task_h  = fl_h//2 - bp - 8
    add_rect(s, tasks_x, fl_y+bp, tasks_w, task_h, fill=PRIMARY)
    add_text(s, 'DIGITAL', tasks_x+8, fl_y+bp+5, tasks_w-16, 10,
             sz=6, bold=True, color=rgb('#ffffff99'))
    add_text(s, 'Digital Tasks → Connections', tasks_x+8, fl_y+bp+16, tasks_w-16, 16,
             sz=9.5, bold=True, color=WHITE)
    add_text(s, 'API retrieves records; sends to HealthSource on success',
             tasks_x+8, fl_y+bp+33, tasks_w-16, 28, sz=7.5, color=rgb('#ffffffb3'), wrap=True)

    add_rect(s, tasks_x, fl_y + fl_h//2 + 2, tasks_w, task_h, fill=BATCH_BG)
    add_text(s, 'MANUAL', tasks_x+8, fl_y + fl_h//2 + 7, tasks_w-16, 10,
             sz=6, bold=True, color=rgb('#ffffff99'))
    add_text(s, 'Group Manual Tasks & Assignment', tasks_x+8, fl_y + fl_h//2 + 18, tasks_w-16, 16,
             sz=9.5, bold=True, color=WHITE)
    add_text(s, 'Onsite specialist or remote ops staff picks up tasks',
             tasks_x+8, fl_y + fl_h//2 + 35, tasks_w-16, 28, sz=7.5, color=rgb('#ffffffb3'), wrap=True)

    mfr_x = tasks_x + tasks_w + 4
    add_text(s, '→', mfr_x, fl_y + fl_h//2 - 8, 16, 16, sz=10, color=DIMMED, align=PP_ALIGN.CENTER)

    # Assembly
    asm_x = mfr_x + 16
    asm_w = bound_x + bound_w - bp - asm_x
    add_rect(s, asm_x, fl_y+bp, asm_w, eng_inner_h, fill=PRIMARY_800)
    add_text(s, 'ASSEMBLY', asm_x+6, fl_y+bp+6, asm_w-12, 10,
             sz=5.5, bold=True, color=rgb('#ffffff66'), align=PP_ALIGN.CENTER)
    add_text(s, 'Merge All\nRetrieved\nDocuments', asm_x+6, fl_y+bp+18, asm_w-12, 46,
             sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
    add_text(s, 'One complete request PDF from all tasks and systems',
             asm_x+6, fl_y+bp+68, asm_w-12, 40, sz=7, color=rgb('#ffffffb3'),
             align=PP_ALIGN.CENTER, wrap=True)

    add_text(s, '→', farr2_x, fl_y + fl_h//2 - 8, farr_w, 16, sz=12, color=DIMMED, align=PP_ALIGN.CENTER)

    # QC
    add_rect(s, qc_x, fl_y, qc_w, fl_h//2 - 4, fill=SURF, lc=BORDER, lw=0.75)
    add_text(s, 'QC', qc_x+4, fl_y+6, qc_w-8, 10, sz=6.5, bold=True, color=DIMMED)
    add_text(s, 'Verify\nAssist', qc_x+4, fl_y+18, qc_w-8, 22, sz=9, color=BODY, wrap=True)
    add_text(s, 'Manual QC', qc_x+4, fl_y+40, qc_w-8, 12, sz=7, color=DIMMED)

    add_text(s, '→', farr3_x, fl_y + fl_h//4 - 8, farr_w, 16, sz=12, color=DIMMED, align=PP_ALIGN.CENTER)

    # Delivery
    add_rect(s, del_x, fl_y, del_w, fl_h//2 - 4, fill=SURF, lc=BORDER, lw=0.75)
    add_text(s, 'DELIVERY', del_x+4, fl_y+6, del_w-8, 10, sz=6.5, bold=True, color=DIMMED)
    add_text(s, 'AutoFax', del_x+4, fl_y+18, del_w-8, 14, sz=9, color=BODY)
    add_text(s, 'Fax, Mail,\nElectronic', del_x+4, fl_y+32, del_w-8, 18, sz=7, color=DIMMED, wrap=True)

    # Foundation strip
    fnd_y = fl_y + fl_h + 16
    fnd_h = 66
    add_rect(s, P, fnd_y, 912, fnd_h, fill=SURF, lc=BORDER, lw=0.75)
    add_text(s, '05 — FOUNDATION LAYER (runs across all above)', P+10, fnd_y+6, 340, 10,
             sz=6.5, bold=True, color=DIMMED)
    fnd_items = [
        ('Audit Trails',
         'Per-task record: system, operations staff, what retrieved, when. Powers compliance and dispute resolution.',
         PRIMARY),
        ('Data Instrumentation',
         'Time-on-task by type, system, and operations staff — enables true labor cost measurement and automation targeting.',
         TEAL),
        ('Exception Handling',
         'When digital retrieval fails, that task reroutes to manual without blocking other tasks in progress.',
         AMBER_BD),
    ]
    fi_w = (912 - 20 - 2*20) // 3
    for fidx, (fname, fdesc, fc) in enumerate(fnd_items):
        fx = P + 10 + fidx * (fi_w + 20)
        add_rect(s, fx, fnd_y+18, 3, fnd_h-22, fill=fc)
        add_text(s, fname,  fx+7, fnd_y+18, fi_w-10, 12, sz=8.5, bold=True, color=BODY)
        add_text(s, fdesc,  fx+7, fnd_y+31, fi_w-10, fnd_h-35, sz=7, color=DIMMED, wrap=True)

    # Registry note
    reg_y = fnd_y + fnd_h + 6
    add_text(s, 'The Task Engine reads from the Record Location Registry — a health system map of which systems hold which record types and date ranges, and which are digitally eligible.',
             P, reg_y, 912, 16, sz=7, italic=True, color=DIMMED, align=PP_ALIGN.CENTER, wrap=True)

    wm(s)
    return s


# ── SLIDE 6: Roadmap ─────────────────────────────────────────────────────────

def build_slide6(prs, layout):
    print('  Slide 6 — Roadmap')
    s = prs.slides.add_slide(layout)
    add_rect(s, 0, 0, 960, 540, fill=WHITE)
    P, H = 26, 20

    add_text(s, 'Roadmap & Resourcing Ask', P, H, 300, 12, sz=7.5, color=PRIMARY)
    add_text(s, 'Validate Before We Scale — Resources Needed by Phase',
             P, H+14, 912, 22, sz=13.5, bold=True, color=BODY)
    add_rect(s, P, H+38, 3, 18, fill=PRIMARY)
    add_text(s, 'Each phase is gated on validating labor cost savings. We do not advance to the next phase without confirming measurable efficiency gains.',
             P+7, H+38, 900, 18, sz=7.5, color=DIMMED, wrap=True)

    # Roadmap
    ask_h = 34
    rm_y = H + 62
    rm_h = 540 - rm_y - ask_h - 14
    gate_w = 44
    phase_w = (912 - 2*gate_w) // 3

    p1x = P
    g1x = p1x + phase_w
    p2x = g1x + gate_w
    g2x = p2x + phase_w
    p3x = g2x + gate_w
    hdr_h = 50

    phases = [
        {
            'x': p1x, 'q': 'Q3 2026', 'label': 'Discovery & Validation',
            'hdr_bg': SURF, 'hdr_txt': BODY, 'lbl_c': DIMMED,
            'dashed': False, 'cond': None,
            'acts': [
                'Stakeholder discovery — operations staff, site managers, health system partners',
                'Prototype task-based fulfillment flow; test with operations staff',
                'Architecture design — Task Engine & HealthSource integration',
                'Define success metrics and validation framework',
            ],
            'chips': [
                ('Product', rgb('#dbeafe'), rgb('#1e40af')),
                ('Design',  rgb('#fce7f3'), rgb('#9d174d')),
                ('Ops',     rgb('#dcfce7'), rgb('#166534')),
            ],
            'chip_lbl_c': DIMMED,
        },
        {
            'x': p2x, 'q': 'Q4 2026 – Q1 2027', 'label': 'MVP Build & Pilot',
            'hdr_bg': PRIMARY, 'hdr_txt': WHITE, 'lbl_c': rgb('#ffffffb3'),
            'dashed': False, 'cond': 'Conditional on Phase 1',
            'acts': [
                'Build Task Engine, routing logic, Task Assignment UI/UX',
                'Integrate Connections (digital) + manual task queue',
                'Live pilot — parallel task routing against real requests',
                'Instrument task-level time tracking; measure cost vs. baseline',
            ],
            'chips': [
                ('Product',           rgb('#dbeafe'), rgb('#1e40af')),
                ('Ops',               rgb('#dcfce7'), rgb('#166534')),
                ('Eng — Connections', DARK_NAVY,     WHITE),
                ('Eng — ROI Zone',    DARK_NAVY,     WHITE),
            ],
            'chip_lbl_c': rgb('#ffffff7f'),
        },
        {
            'x': p3x, 'q': 'Q2 2027+', 'label': 'Scale to More Sites',
            'hdr_bg': PRIMARY_LT, 'hdr_txt': PRIMARY_800, 'lbl_c': PRIMARY,
            'dashed': True, 'cond': 'Conditional on Phase 2',
            'acts': [
                'Onboard additional health systems & site configurations',
                'Scale task routing infrastructure from pilot learnings',
                'Refine operations staff UX; measure aggregate cost savings',
                'Build business case for full network rollout',
            ],
            'chips': [
                ('Product',           rgb('#dbeafe'), rgb('#1e40af')),
                ('Ops',               rgb('#dcfce7'), rgb('#166534')),
                ('Eng — Connections', DARK_NAVY,     WHITE),
                ('Eng — ROI Zone',    DARK_NAVY,     WHITE),
            ],
            'chip_lbl_c': DIMMED,
        },
    ]

    for ph in phases:
        px2 = ph['x']
        if ph['dashed']:
            add_dashed(s, px2, rm_y, phase_w, rm_h, lc=rgb('#2945f059'), fill=WHITE, lw=1.5)
        else:
            lc = PRIMARY if ph['hdr_bg'] == PRIMARY else BORDER
            add_rect(s, px2, rm_y, phase_w, rm_h, fill=WHITE, lc=lc, lw=1)

        offset = 0
        if ph['cond']:
            add_rect(s, px2+5, rm_y+5, phase_w-10, 14, fill=AMBER_LT, lc=rgb('#f59e0b66'), lw=0.5)
            add_text(s, '⚑ ' + ph['cond'], px2+8, rm_y+7, phase_w-16, 11,
                     sz=6, bold=True, color=AMBER)
            offset = 20

        add_rect(s, px2, rm_y+offset, phase_w, hdr_h-offset, fill=ph['hdr_bg'])
        add_text(s, ph['q'], px2+10, rm_y+offset+6, phase_w-20, 16,
                 sz=11, bold=True, color=ph['hdr_txt'])
        add_text(s, ph['label'], px2+10, rm_y+offset+24, phase_w-20, 12,
                 sz=7.5, bold=True, color=ph['lbl_c'])

        body_y = rm_y + hdr_h + 5
        add_text(s, 'SCOPE OF WORK', px2+10, body_y, phase_w-20, 10,
                 sz=6.5, bold=True, color=DIMMED)
        body_y += 13
        for act in ph['acts']:
            add_rect(s, px2+10, body_y+5, 4, 4, fill=BODY)
            add_text(s, act, px2+18, body_y, phase_w-28, 20, sz=7.5, color=BODY, wrap=True)
            body_y += 22

        res_y = rm_y + rm_h - 46
        border_c = rgb('#ffffff30') if ph['hdr_bg'] == PRIMARY else BORDER
        add_rect(s, px2+6, res_y, phase_w-12, 1, fill=border_c)
        add_text(s, 'RESOURCES NEEDED', px2+10, res_y+4, phase_w-20, 10,
                 sz=6.5, bold=True, color=ph['chip_lbl_c'])
        chip_x = px2 + 10
        chip_y = res_y + 16
        for (clbl, cbg, ctc) in ph['chips']:
            cw = max(38, len(clbl)*5+12)
            add_rect(s, chip_x, chip_y, cw, 14, fill=cbg)
            add_text(s, clbl, chip_x+3, chip_y+2, cw-6, 10, sz=6.5, bold=True, color=ctc)
            chip_x += cw + 4
            if chip_x > px2 + phase_w - 20:
                chip_x = px2 + 10
                chip_y += 17

    # Gate connectors
    for gx in [g1x, g2x]:
        g_cx = gx + gate_w//2
        g_cy = rm_y + rm_h//2
        add_rect(s, g_cx-1, rm_y+8, 2, rm_h//2 - 20, fill=AMBER_BD)
        add_rect(s, g_cx-1, g_cy+18, 2, rm_h//2 - 24, fill=AMBER_BD)
        diamond = add_rect(s, g_cx-12, g_cy-12, 24, 24, fill=AMBER)
        diamond.rotation = 45
        add_text(s, '✓', g_cx-7, g_cy-7, 14, 14, sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, 'If\nvalidated', gx, g_cy+20, gate_w, 18,
                 sz=6, bold=True, color=AMBER, align=PP_ALIGN.CENTER, wrap=True)

    # Ask bar
    ask_y = rm_y + rm_h + 6
    add_rect(s, P, ask_y, 912, ask_h, fill=PRIMARY_LT, lc=PRIMARY, lw=0.75)
    add_text(s, 'IMMEDIATE ASK', P+10, ask_y+8, 110, 10, sz=7.5, bold=True, color=PRIMARY_800)
    add_rect(s, P+118, ask_y+7, 1, ask_h-14, fill=rgb('#2945f033'))
    asks = [
        ('Design', '1 resource for Q3 discovery & prototype'),
        ('Ops', 'embedded partner Q3 through Q4'),
        ('Eng (Connections + ROI Zone)', 'scoping & commitment for Q4'),
    ]
    ai_x = P + 126
    for idx, (albl, adesc) in enumerate(asks):
        add_text(s, albl, ai_x, ask_y+5, max(40, len(albl)*5+10), 12,
                 sz=8, bold=True, color=PRIMARY_800)
        add_text(s, '— ' + adesc, ai_x + max(40, len(albl)*5+10) + 2, ask_y+5,
                 220, 12, sz=8, color=PRIMARY_800)
        if idx < 2:
            sep_x = ai_x + max(40, len(albl)*5+10) + 220 + 8
            add_text(s, '·', sep_x, ask_y+5, 12, 12, sz=10, color=rgb('#2945f066'),
                     align=PP_ALIGN.CENTER)
        ai_x += max(40, len(albl)*5+10) + 220 + 22

    wm(s)
    return s


# ── SLIDE 2: Executive Summary ───────────────────────────────────────────────

def build_slide_exec(prs, layout):
    print('  Slide 2 — Executive Summary')
    s = prs.slides.add_slide(layout)
    add_rect(s, 0, 0, 960, 540, fill=WHITE)
    P, H = 52, 44

    add_text(s, 'Executive Summary', P, H, 300, 12, sz=7.5, color=PRIMARY)
    add_text(s, 'Each ROI request is multiple independent tasks across different systems and locations. Completing them sequentially creates unnecessary delay, repeat work, and labor costs we can\'t measure.',
             P, H+20, 820, 88, sz=13, color=BODY, wrap=True)

    # Three callouts
    co_y = H + 145
    co_h = 540 - co_y - 70  # leave room for ask bar
    co_w = (912 - 2*2) // 3  # 2px dividers between callouts
    callouts = [
        ('1', 'One request, many independent tasks', [
            ('Each request spans', 'multiple systems', '— each requiring separate login, credentials, and context'),
            ('Some tasks require', 'physical retrieval', '— leaving the desk, traveling onsite or coordinating off-site'),
            ('Today,', 'one person handles all of them', 'in sequence, one at a time'),
        ]),
        ('2', 'Digital and manual worked sequentially', [
            ('Requests have', 'digital components', '(EHR records via Connections) and manual components (onsite imaging, archival)'),
            ('Manual tasks', "can't begin until digital is done", '— or vice versa'),
            ('Sequential execution creates', 'repeat work, delay, and no efficiency gain', ''),
        ]),
        ('3', 'Our hypothesis: complete fulfillment tasks in parallel', [
            ('Split each request into task units — route', 'digital and manual simultaneously', 'to the right processor'),
            ('', 'Faster TAT', '— total time equals the slowest task, not the sum of all'),
            ('', 'Labor visibility', '— task-level data reveals where manual effort goes and what to automate next'),
        ]),
    ]
    for cidx, (num, title, bullets) in enumerate(callouts):
        cx = P + cidx * (co_w + 2)
        # Numbered circle
        add_rect(s, cx, co_y, 28, 28, fill=PRIMARY)
        add_text(s, num, cx, co_y, 28, 28, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_text(s, title, cx+38, co_y+2, co_w-38-32, 22, sz=11, bold=True, color=BODY, wrap=True)
        # Bullets
        by = co_y + 40
        for (pre, strong, post) in bullets:
            add_rect(s, cx+38, by+5, 5, 5, fill=PRIMARY)
            full = (pre + ' ' if pre else '') + strong + (' ' + post if post else '')
            add_text(s, full.strip(), cx+48, by, co_w-48-32, 28, sz=10.5, color=DIMMED, wrap=True)
            by += 32
        # Divider (right edge, except last)
        if cidx < 2:
            add_rect(s, cx + co_w, co_y, 2, co_h, fill=BORDER)

    # Ask bar
    ask_y = 540 - 52
    add_rect(s, P, ask_y, 856, 38, fill=PRIMARY)
    add_text(s, 'THE ASK', P+14, ask_y+10, 70, 12, sz=8, bold=True, color=rgb('#ffffff99'))
    add_rect(s, P+84, ask_y+9, 1, 20, fill=rgb('#ffffff40'))
    add_text(s, 'Design + Ops to prototype and validate in Q3  ·  Eng scoping (Connections + ROI Zone) commitment for Q4 if the hypothesis is confirmed',
             P+96, ask_y+9, 756, 20, sz=10, color=WHITE, wrap=True)

    wm(s)
    return s


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    build_slide1(prs, blank)
    build_slide_exec(prs, blank)
    build_slide2(prs, blank)
    build_slide3(prs, blank)
    build_slide4(prs, blank)
    build_slide5(prs, blank)
    build_slide6(prs, blank)

    prs.save(str(OUT))
    print(f'\nSaved → {OUT}')
    print('All shapes are native PPTX — every element is editable in Google Slides or PowerPoint.')

main()
